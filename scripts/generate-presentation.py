# /// script
# requires-python = ">=3.10"
# dependencies = ["python-pptx>=1.0"]
# ///
"""
Genera la presentación PPTX de producto ClimaLab.
Uso: uv run scripts/generate-presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Brand colors ──
GREEN = RGBColor(0x28, 0x94, 0x48)
CYAN = RGBColor(0x1F, 0xAC, 0xC0)
RED = RGBColor(0xC3, 0x24, 0x21)
DARK = RGBColor(0x15, 0x14, 0x14)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF4, 0xF7, 0xF5)
MEDIUM_GRAY = RGBColor(0x6B, 0x6B, 0x6B)
LIGHT_GRAY = RGBColor(0xE5, 0xE5, 0xE5)
TABLE_HEADER_BG = RGBColor(0x22, 0x7A, 0x3C)
TABLE_ALT_BG = RGBColor(0xF0, 0xF7, 0xF2)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_rect(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=14,
                    color=DARK, bullet_color=GREEN, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
        # Bullet character
        pPr = p._pPr
        if pPr is None:
            from pptx.oxml.ns import qn
            pPr = p._p.get_or_add_pPr()
        from pptx.oxml.ns import qn
        buFont = pPr.makeelement(qn('a:buFont'), {'typeface': 'Arial'})
        buChar = pPr.makeelement(qn('a:buChar'), {'char': '●'})
        # Remove existing bullets
        for child in list(pPr):
            if child.tag.endswith('buChar') or child.tag.endswith('buFont') or child.tag.endswith('buNone'):
                pPr.remove(child)
        pPr.append(buFont)
        pPr.append(buChar)
    return txBox


def add_table(slide, left, top, width, height, rows_data, col_widths=None):
    """rows_data: list of lists. First row = header."""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r, row in enumerate(rows_data):
        for c, cell_text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(cell_text)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.name = "Calibri"
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                    paragraph.alignment = PP_ALIGN.CENTER
                else:
                    paragraph.font.color.rgb = DARK

            # Cell fill
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER_BG
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT_BG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE

    return table_shape


def make_section_header(slide, number, title, subtitle=None):
    """Green bar at top with section number + title."""
    set_slide_bg(slide, WHITE)
    # Top accent bar
    add_shape_rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.06), GREEN)
    # Section number circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(0.6), Inches(0.4), Inches(0.55), Inches(0.55)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = GREEN
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(2)

    # Title
    add_text_box(slide, Inches(1.35), Inches(0.38), Inches(10), Inches(0.55),
                 title, font_size=26, color=DARK, bold=True)
    if subtitle:
        add_text_box(slide, Inches(1.35), Inches(0.88), Inches(10), Inches(0.4),
                     subtitle, font_size=14, color=MEDIUM_GRAY)


# ══════════════════════════════════════════════════════════════════════
#  SLIDES
# ══════════════════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]  # blank layout

# ── SLIDE 1: Portada ──
sl = prs.slides.add_slide(BLANK)
set_slide_bg(sl, GREEN)
# Subtle accent bar
add_shape_rect(sl, Inches(0), Inches(0), Inches(13.33), Inches(0.12), RGBColor(0x1E, 0x7A, 0x38))
# Logo placeholder
add_text_box(sl, Inches(0.8), Inches(0.5), Inches(2), Inches(0.5),
             "ClimaLab", font_size=22, color=WHITE, bold=True)
# Main title
add_text_box(sl, Inches(0.8), Inches(2.2), Inches(11), Inches(1.2),
             "Mide lo que importa\nen tu organización",
             font_size=48, color=WHITE, bold=True)
# Subtitle
add_text_box(sl, Inches(0.8), Inches(3.8), Inches(9), Inches(0.8),
             "Plataforma de diagnóstico de clima organizacional que transforma\ndatos en decisiones accionables sobre tu gente.",
             font_size=20, color=RGBColor(0xD0, 0xF0, 0xD8))
# Bottom bar
add_shape_rect(sl, Inches(0), Inches(6.8), Inches(13.33), Inches(0.7), RGBColor(0x1E, 0x7A, 0x38))
add_text_box(sl, Inches(0.8), Inches(6.88), Inches(5), Inches(0.45),
             "Un producto de Rizo.ma Consulting  ·  Panamá",
             font_size=14, color=RGBColor(0xA0, 0xD8, 0xB0))
add_text_box(sl, Inches(8), Inches(6.88), Inches(4.5), Inches(0.45),
             "climalab.vercel.app",
             font_size=14, color=RGBColor(0xA0, 0xD8, 0xB0), alignment=PP_ALIGN.RIGHT)


# ── SLIDE 2: El problema ──
sl = prs.slides.add_slide(BLANK)
make_section_header(sl, 1, "El problema",
                    "Las organizaciones toman decisiones sobre su gente basadas en intuición")

problems = [
    "Las encuestas genéricas producen datos sin contexto ni rigor científico",
    "Los resultados llegan tarde, sin segmentación y sin recomendaciones claras",
    "Contratar consultoría tradicional es costoso y lento para PYMEs",
    "Falta de anonimato real reduce la participación y honestidad",
    "Sin análisis avanzado, las intervenciones se hacen a ciegas",
]
add_bullet_list(sl, Inches(0.8), Inches(1.6), Inches(7), Inches(3.5),
                problems, font_size=16, spacing=Pt(12))

# Right side callout
callout = add_shape_rect(sl, Inches(8.5), Inches(1.8), Inches(4.2), Inches(2.2), RGBColor(0xFE, 0xF2, 0xF2))
add_text_box(sl, Inches(8.8), Inches(2.0), Inches(3.6), Inches(1.8),
             "El resultado:\n\nIntervenciones a ciegas, rotación evitable y talento que se pierde.",
             font_size=16, color=RED, bold=True)

# Bottom insight
add_shape_rect(sl, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.2), LIGHT_BG)
add_text_box(sl, Inches(1.2), Inches(5.7), Inches(11), Inches(0.8),
             "ClimaLab resuelve esto: rigor de consultoría + velocidad de plataforma + IA integrada",
             font_size=18, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)


# ── SLIDE 3: La solución — comparativa ──
sl = prs.slides.add_slide(BLANK)
make_section_header(sl, 2, "La solución",
                    "Rigor de consultoría, velocidad de plataforma, costo accesible")

table_data = [
    ["Criterio", "Encuestas genéricas", "Consultoría tradicional", "ClimaLab"],
    ["Fundamento científico", "Bajo", "Alto", "Alto (22 dims validadas)"],
    ["Velocidad de resultados", "Media", "Lenta (semanas)", "Inmediata"],
    ["Costo", "Bajo", "Alto", "Accesible"],
    ["Análisis estadístico", "Descriptivo básico", "Manual (consultor)", "Automatizado (rwg, α, Pearson)"],
    ["Inteligencia Artificial", "No", "No", "Sí (triple backend)"],
    ["Análisis de redes (ONA)", "No", "Raro", "Sí (Leiden + NMI)"],
    ["Marca blanca", "Limitada", "Variable", "Completa"],
    ["Contexto LATAM", "Genérico", "Depende del consultor", "Nativo"],
]
add_table(sl, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.2), table_data,
          col_widths=[Inches(2.4), Inches(2.8), Inches(3.0), Inches(3.5)])


# ── SLIDE 4: Propuesta de valor ──
sl = prs.slides.add_slide(BLANK)
make_section_header(sl, 3, "Propuesta de valor", "Mide · Entiende · Actúa")

# Three columns
cols = [
    ("Mide", "Diagnóstico de 22 dimensiones científicamente validadas en menos de 10 minutos por persona.",
     "📊", GREEN),
    ("Entiende", "Visualiza resultados por área, antigüedad, género y perfiles de compromiso de forma inmediata.",
     "👁", CYAN),
    ("Actúa", "Recibe recomendaciones basadas en evidencia e impulsadas por IA para intervenir donde más importa.",
     "🎯", RGBColor(0xFF, 0x80, 0x44)),
]

for i, (title, desc, icon, color) in enumerate(cols):
    x = Inches(0.8 + i * 4.1)
    # Card background
    card = add_shape_rect(sl, x, Inches(1.8), Inches(3.7), Inches(4.5), LIGHT_BG)
    card.shadow.inherit = False
    # Color accent top
    add_shape_rect(sl, x, Inches(1.8), Inches(3.7), Inches(0.08), color)
    # Icon
    add_text_box(sl, x + Inches(0.3), Inches(2.1), Inches(1), Inches(0.7),
                 icon, font_size=36, color=color)
    # Title
    add_text_box(sl, x + Inches(0.3), Inches(2.8), Inches(3.1), Inches(0.5),
                 title, font_size=24, color=DARK, bold=True)
    # Description
    add_text_box(sl, x + Inches(0.3), Inches(3.4), Inches(3.1), Inches(2.5),
                 desc, font_size=15, color=MEDIUM_GRAY)


# ── SLIDE 5: Instrumento — 22 dimensiones ──
sl = prs.slides.add_slide(BLANK)
make_section_header(sl, 4, "Instrumento científico — ClimaLab Core v4.0",
                    "22 dimensiones en 4 categorías + Engagement transversal")

categories = [
    ("Bienestar", "6 dimensiones",
     "Orgullo Institucional · Propósito del Trabajo · Seguridad Física y Psicológica · Balance Vida-Trabajo · Cuidado Mutuo · Demandas Laborales",
     GREEN),
    ("Dirección y Supervisión", "5 dimensiones",
     "Liderazgo Efectivo · Autonomía · Comunicación Interna · Confianza Institucional · Claridad de Rol",
     CYAN),
    ("Compensación", "5 dimensiones",
     "Compensación · Reconocimiento · Beneficios · Equidad en Ascensos · No Discriminación e Inclusión",
     RGBColor(0xFF, 0x80, 0x44)),
    ("Cultura", "5 dimensiones",
     "Cohesión de Equipo · Innovación y Cambio · Resultados y Logros · Desarrollo Profesional · Aprendizaje Organizacional",
     RGBColor(0x2F, 0x5D, 0xFF)),
]

for i, (cat_name, count, dims, color) in enumerate(categories):
    row = i // 2
    col = i % 2
    x = Inches(0.8 + col * 6.2)
    y = Inches(1.5 + row * 2.6)

    add_shape_rect(sl, x, y, Inches(5.8), Inches(2.3), WHITE)
    add_shape_rect(sl, x, y, Inches(0.08), Inches(2.3), color)  # Left accent
    add_text_box(sl, x + Inches(0.3), y + Inches(0.15), Inches(4), Inches(0.4),
                 cat_name, font_size=18, color=color, bold=True)
    add_text_box(sl, x + Inches(4.2), y + Inches(0.15), Inches(1.3), Inches(0.4),
                 count, font_size=12, color=MEDIUM_GRAY, alignment=PP_ALIGN.RIGHT)
    add_text_box(sl, x + Inches(0.3), y + Inches(0.65), Inches(5.2), Inches(1.5),
                 dims, font_size=13, color=DARK)

# Engagement bar at bottom
add_shape_rect(sl, Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.5), GREEN)
add_text_box(sl, Inches(1.2), Inches(6.72), Inches(10), Inches(0.45),
             "Variable transversal: Engagement y Compromiso (UWES-9) — 5 ítems como variable dependiente",
             font_size=13, color=WHITE, bold=True)


# ── SLIDE 6: Modalidades ──
sl = prs.slides.add_slide(BLANK)
make_section_header(sl, 5, "Modalidades de medición",
                    "Instrumento flexible según la necesidad")

table_data = [
    ["Modalidad", "Ítems", "Uso ideal", "Duración"],
    ["Core v4.0", "109 (107 + 2 checks)", "Diagnóstico integral anual o semestral", "~10 min"],
    ["Pulso v4.0", "22 ítems ancla", "Pulsos rápidos trimestrales o mensuales", "~3 min"],
]
add_table(sl, Inches(0.8), Inches(1.6), Inches(11.7), Inches(1.8), table_data,
          col_widths=[Inches(2.5), Inches(3), Inches(4.2), Inches(2)])

add_text_box(sl, Inches(0.8), Inches(3.8), Inches(5), Inches(0.4),
             "Módulos opcionales (se agregan al instrumento base)",
             font_size=18, color=DARK, bold=True)

mod_data = [
    ["Módulo", "Ítems", "Aplicación"],
    ["Gestión del Cambio (CAM)", "8", "Procesos de transformación organizacional"],
    ["Orientación al Cliente (CLI)", "4", "Fortalecimiento de cultura de servicio"],
    ["Preparación Digital (DIG)", "4", "Evaluación de madurez tecnológica"],
]
add_table(sl, Inches(0.8), Inches(4.3), Inches(11.7), Inches(2.2), mod_data,
          col_widths=[Inches(3.5), Inches(1.5), Inches(6.7)])

add_text_box(sl, Inches(0.8), Inches(6.7), Inches(11), Inches(0.4),
             "Cada dimensión respaldada por teoría peer-reviewed: Mael & Ashforth, Deci & Ryan, Eisenberger, Edmondson, Locke & Latham, etc.",
             font_size=12, color=MEDIUM_GRAY)


# ── SLIDE 7: Experiencia de encuesta ──
sl = prs.slides.add_slide(BLANK)
make_section_header(sl, 6, "Experiencia de encuesta",
                    "Diseñada para maximizar participación y calidad de datos")

steps = [
    ("1", "Bienvenida", "Personalizada con logo\ny colores de la org"),
    ("2", "Demografía", "Departamento,\nantigüedad, género"),
    ("3", "Dimensiones", "Ítems aleatorizados\nen escala Likert 1-5"),
    ("4", "Abiertas + eNPS", "Fortalezas, mejoras,\ncomentarios + eNPS 0-10"),
    ("5", "Agradecimiento", "Mensaje personalizado\ncon branding"),
]

for i, (num, title, desc) in enumerate(steps):
    x = Inches(0.5 + i * 2.5)
    # Circle with number
    circle = sl.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.7), Inches(1.7), Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = GREEN
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    # Arrow (except last)
    if i < 4:
        add_text_box(sl, x + Inches(2.0), Inches(1.75), Inches(0.5), Inches(0.5),
                     "→", font_size=24, color=LIGHT_GRAY, bold=True)

    add_text_box(sl, x + Inches(0.15), Inches(2.5), Inches(2.2), Inches(0.35),
                 title, font_size=15, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.15), Inches(2.9), Inches(2.2), Inches(0.8),
                 desc, font_size=12, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

# Guarantees section
add_text_box(sl, Inches(0.8), Inches(4.2), Inches(5), Inches(0.4),
             "Garantías de calidad y anonimato", font_size=18, color=DARK, bold=True)

guarantees = [
    "100% anónima — tokens únicos, sin email ni login requerido",
    "Respaldo automático en localStorage (recuperación ante cierre)",
    "Responsive — optimizada para móvil y escritorio",
    "Marca blanca completa (colores, logo, textos personalizados)",
    "Checks de atención para filtrar respuestas sin calidad",
    "Ítems invertidos para detectar sesgo de aquiescencia",
]
add_bullet_list(sl, Inches(0.8), Inches(4.7), Inches(11), Inches(2.5),
                guarantees, font_size=14, spacing=Pt(8))


# ── SLIDE 8: Dashboard de resultados ──
sl = prs.slides.add_slide(BLANK)
make_section_header(sl, 7, "Dashboard de resultados — 11 módulos",
                    "Análisis completo en un solo lugar")

modules = [
    ["#", "Módulo", "Qué muestra"],
    ["1", "Dashboard", "KPIs, perfiles de engagement, eNPS, fortalezas y debilidades"],
    ["2", "Dimensiones", "Puntajes detallados, rankings, heatmaps por segmento"],
    ["3", "Tendencias", "Evolución histórica entre campañas"],
    ["4", "Segmentos", "Análisis por departamento, antigüedad y género"],
    ["5", "Benchmarks", "Comparativas internas entre áreas"],
    ["6", "Drivers", "Correlaciones Pearson, impulsores del engagement"],
    ["7", "Alertas", "Detección automática de crisis y grupos de riesgo"],
    ["8", "Comentarios", "Análisis de sentimiento y extracción de temas"],
    ["9", "Red Perceptual", "Análisis de redes organizacionales (ONA)"],
    ["10", "Ficha Técnica", "Metodología, confiabilidad, limitaciones auto-detectadas"],
    ["11", "Exportar", "DOCX editable, Excel, CSV, reporte IA"],
]
add_table(sl, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.5), modules,
          col_widths=[Inches(0.7), Inches(2.2), Inches(8.8)])


# ── SLIDE 9: Rigor estadístico ──
sl = prs.slides.add_slide(BLANK)
make_section_header(sl, 8, "Rigor estadístico",
                    "Métodos psicométricos de nivel académico, completamente automatizados")

stats_data = [
    ["Método", "Propósito", "Umbral"],
    ["Cronbach's Alpha", "Consistencia interna por dimensión", "≥0.70 aceptable"],
    ["rwg(j) (James et al.)", "Acuerdo intra-grupo por dimensión", "≥0.70 suficiente"],
    ["Correlación de Pearson", "Identificar drivers del engagement", "n ≥ 10"],
    ["eNPS (0-10)", "Lealtad y recomendación", "Promotores ≥9, Detractores ≤6"],
    ["Margen de error", "Precisión muestral con corrección FPC", "IC 95%"],
    ["Favorabilidad", "% de respuestas ≥4 en escala de 5 puntos", "—"],
]
add_table(sl, Inches(0.8), Inches(1.5), Inches(11.7), Inches(3.6), stats_data,
          col_widths=[Inches(3), Inches(5.7), Inches(3)])

# Engagement profiles
add_text_box(sl, Inches(0.8), Inches(5.4), Inches(5), Inches(0.4),
             "Perfiles de Engagement:", font_size=16, color=DARK, bold=True)

profiles = [
    ("Embajadores", "≥ 4.5", GREEN),
    ("Comprometidos", "4.0 – 4.49", CYAN),
    ("Neutrales", "3.0 – 3.99", RGBColor(0xFF, 0x80, 0x44)),
    ("Desvinculados", "< 3.0", RED),
]
for i, (name, threshold, color) in enumerate(profiles):
    x = Inches(0.8 + i * 3.0)
    add_shape_rect(sl, x, Inches(5.9), Inches(2.6), Inches(0.8), color)
    add_text_box(sl, x + Inches(0.2), Inches(5.95), Inches(2.2), Inches(0.35),
                 name, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.2), Inches(6.3), Inches(2.2), Inches(0.35),
                 threshold, font_size=12, color=WHITE, alignment=PP_ALIGN.CENTER)

add_text_box(sl, Inches(0.8), Inches(6.9), Inches(11), Inches(0.4),
             "Detección automática de limitaciones: alpha bajo, rwg bajo, tasa de respuesta baja, muestra pequeña. Segmentos <5 personas excluidos.",
             font_size=11, color=MEDIUM_GRAY)


# ── SLIDE 10: ONA ──
sl = prs.slides.add_slide(BLANK)
make_section_header(sl, 9, "Análisis de Redes Organizacionales (ONA)",
                    "Más allá de las encuestas: mapea la percepción compartida")

# Left column
add_text_box(sl, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.4),
             "¿Qué es?", font_size=18, color=DARK, bold=True)
add_text_box(sl, Inches(0.8), Inches(2.0), Inches(5.5), Inches(1.2),
             "Análisis de redes perceptuales que identifica cómo se agrupan los colaboradores por su visión compartida de la organización.\n\nNO es ONA sociométrica — mide percepción, no interacción.",
             font_size=14, color=DARK)

add_text_box(sl, Inches(0.8), Inches(3.4), Inches(5.5), Inches(0.4),
             "Metodología", font_size=18, color=DARK, bold=True)
method_items = [
    "Similitud coseno entre vectores de 22 dimensiones",
    "Algoritmo Leiden (50 iteraciones con estabilidad NMI)",
    "Estabilidad: >0.80 robusto, 0.50-0.80 moderado",
    "Motor: python-igraph (kernel C, alta performance)",
]
add_bullet_list(sl, Inches(0.8), Inches(3.8), Inches(5.5), Inches(2.0),
                method_items, font_size=13, spacing=Pt(8))

# Right column - metrics
add_shape_rect(sl, Inches(7), Inches(1.5), Inches(5.8), Inches(5.0), LIGHT_BG)
add_text_box(sl, Inches(7.3), Inches(1.7), Inches(5.2), Inches(0.4),
             "Métricas y visualizaciones", font_size=18, color=GREEN, bold=True)

ona_metrics = [
    "Comunidades perceptuales (grupos con visión similar)",
    "Nodos puente — traductores culturales entre áreas",
    "Bordes críticos — conexiones vulnerables",
    "Centralidad: eigenvector, betweenness, grado",
    "Matriz de densidad departamental",
    "Dimensiones discriminantes por comunidad",
    "Grafo PNG generado en servidor (FR layout)",
    "Índice de estabilidad global (NMI)",
]
add_bullet_list(sl, Inches(7.3), Inches(2.2), Inches(5.2), Inches(3.5),
                ona_metrics, font_size=13, spacing=Pt(8))

add_text_box(sl, Inches(0.8), Inches(6.7), Inches(11), Inches(0.4),
             "Mínimo 10 respondientes. Resultados almacenados como JSONB con imagen base64 en campaign_analytics.",
             font_size=11, color=MEDIUM_GRAY)


# ── SLIDE 11: IA integrada ──
sl = prs.slides.add_slide(BLANK)
make_section_header(sl, 10, "Inteligencia Artificial integrada",
                    "Triple backend con fallback automático — narrativas en español")

# Backend architecture
backends = [
    ("Prioridad 1", "Anthropic (Claude Haiku 4.5)", "~2-5s · ~$0.03/análisis", GREEN),
    ("Prioridad 2", "DGX / OpenAI-compatible", "Qwen 2.5 72B vía Cloudflare Tunnel", CYAN),
    ("Prioridad 3", "Ollama nativo", "Auto-hospedado, sin costo", MEDIUM_GRAY),
]
for i, (priority, name, detail, color) in enumerate(backends):
    x = Inches(0.8 + i * 4.1)
    add_shape_rect(sl, x, Inches(1.5), Inches(3.7), Inches(1.5), WHITE)
    add_shape_rect(sl, x, Inches(1.5), Inches(3.7), Inches(0.06), color)
    add_text_box(sl, x + Inches(0.2), Inches(1.65), Inches(3.3), Inches(0.3),
                 priority, font_size=11, color=color, bold=True)
    add_text_box(sl, x + Inches(0.2), Inches(1.95), Inches(3.3), Inches(0.35),
                 name, font_size=16, color=DARK, bold=True)
    add_text_box(sl, x + Inches(0.2), Inches(2.4), Inches(3.3), Inches(0.4),
                 detail, font_size=12, color=MEDIUM_GRAY)

    # Arrow between cards
    if i < 2:
        add_text_box(sl, x + Inches(3.5), Inches(1.9), Inches(0.8), Inches(0.5),
                     "→", font_size=28, color=LIGHT_GRAY, bold=True, alignment=PP_ALIGN.CENTER)

# Analysis types table
ai_data = [
    ["Análisis", "Página", "Genera"],
    ["Narrativa ejecutiva", "Dashboard", "Resumen, hallazgos, preocupaciones, recomendación"],
    ["Insights de drivers", "Drivers", "Interpretación, paradojas, quick wins"],
    ["Contexto de alertas", "Alertas", "Hipótesis de causa raíz + recomendación"],
    ["Perfiles de segmento", "Segmentos", "Narrativa con fortalezas/riesgos por grupo"],
    ["Análisis de comentarios", "Comentarios", "Temas, sentimiento, resumen por tipo de pregunta"],
    ["Narrativa de tendencias", "Tendencias", "Trayectoria, dimensiones en alza/baja, inflexiones"],
]
add_table(sl, Inches(0.8), Inches(3.3), Inches(11.7), Inches(3.5), ai_data,
          col_widths=[Inches(3), Inches(1.8), Inches(6.9)])

add_text_box(sl, Inches(0.8), Inches(7.0), Inches(11), Inches(0.3),
             "La IA solo resume hallazgos estadísticos existentes — sin riesgo de alucinación. Prompts en español. Fail-fast si no hay proveedor configurado.",
             font_size=11, color=MEDIUM_GRAY)


# ── SLIDE 12: Marca blanca ──
sl = prs.slides.add_slide(BLANK)
make_section_header(sl, 11, "Sistema de marca blanca",
                    "Cada organización ve su propia identidad en todos los puntos de contacto")

# Elements
add_text_box(sl, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.4),
             "Elementos personalizables", font_size=18, color=DARK, bold=True)
brand_items = [
    "Logo (subida y almacenamiento en CDN — Supabase Storage)",
    "Colores: primario, secundario, acento, texto, fondo",
    "Texto de bienvenida y agradecimiento personalizado",
    "Pie de email personalizado",
    'Opción "Powered by ClimaLab" configurable',
    "Editor visual con previsualización en tiempo real",
]
add_bullet_list(sl, Inches(0.8), Inches(2.0), Inches(5.5), Inches(3.0),
                brand_items, font_size=14, spacing=Pt(8))

# Touchpoints
add_shape_rect(sl, Inches(7), Inches(1.5), Inches(5.8), Inches(4.5), LIGHT_BG)
add_text_box(sl, Inches(7.3), Inches(1.7), Inches(5.2), Inches(0.4),
             "Puntos de contacto branded", font_size=18, color=GREEN, bold=True)

touchpoints = [
    ["Punto de contacto", "Personalización"],
    ["Encuesta", "Header, botones, progreso, textos"],
    ["Emails (4 tipos)", "Header con logo, CTA color acento, footer"],
    ["Reporte DOCX", "Portada, encabezados, KPIs, logo embebido"],
    ["Panel de resultados", "Logo en sidebar de navegación"],
]
add_table(sl, Inches(7.2), Inches(2.3), Inches(5.4), Inches(3.0), touchpoints,
          col_widths=[Inches(2.2), Inches(3.2)])


# ── SLIDE 13: Emails ──
sl = prs.slides.add_slide(BLANK)
make_section_header(sl, 12, "Infraestructura de email",
                    "4 tipos de comunicación automatizada y personalizada vía Resend")

email_data = [
    ["Tipo", "Momento", "Contenido"],
    ["Invitación", "Lanzamiento de campaña", "Link de encuesta + contexto organizacional"],
    ["Recordatorio", "Durante campaña activa", "Seguimiento a no completados (con tracking)"],
    ["Campaña cerrada", "Al cerrar medición", "Aviso de fin de período de respuesta"],
    ["Resultados listos", "Post-análisis", "Compartir hallazgos con stakeholders"],
]
add_table(sl, Inches(0.8), Inches(1.5), Inches(11.7), Inches(2.8), email_data,
          col_widths=[Inches(2.5), Inches(3.5), Inches(5.7)])

features = [
    "Layout HTML compartido con logo/colores de la organización",
    "Botón CTA con color acento del brand de la organización",
    "Footer condicional 'Powered by ClimaLab'",
    "Tracking de recordatorios: fecha de último envío + contador por participante",
    "Envío masivo de recordatorios con un clic desde la página de campaña",
]
add_bullet_list(sl, Inches(0.8), Inches(4.8), Inches(11), Inches(2.0),
                features, font_size=14, spacing=Pt(8))


# ── SLIDE 14: Exportación ──
sl = prs.slides.add_slide(BLANK)
make_section_header(sl, 13, "Exportación y reportes",
                    "Lleva los resultados donde los necesites")

export_data = [
    ["Formato", "Contenido", "Uso"],
    ["DOCX (Word)", "Reporte ejecutivo — 14 secciones con marca y logo", "Presentar a directiva, editar"],
    ["Excel", "8 hojas (dims, ítems, segmentos, drivers, alertas, comentarios, ficha)", "Análisis profundo"],
    ["CSV / JSON", "Datos crudos de dimensiones y resultados", "Integración con BI tools"],
    ["Reporte IA", "Narrativa ejecutiva con todos los insights", "Comunicación rápida"],
]
add_table(sl, Inches(0.8), Inches(1.5), Inches(11.7), Inches(2.8), export_data,
          col_widths=[Inches(2), Inches(5.5), Inches(4.2)])

add_text_box(sl, Inches(0.8), Inches(4.6), Inches(5), Inches(0.4),
             "14 secciones del reporte DOCX:", font_size=16, color=DARK, bold=True)

sections = [
    "Portada", "Resumen ejecutivo", "KPIs clave", "Categorías",
    "Dimensiones", "Departamentos", "Alertas (+IA)", "Drivers (+IA)",
    "Comentarios", "Perfiles de segmento", "Tendencias",
    "Indicadores de negocio", "Red perceptual (ONA)", "Ficha técnica"
]
# Two rows of 7 badges
for i, section in enumerate(sections):
    row = i // 7
    col = i % 7
    x = Inches(0.8 + col * 1.75)
    y = Inches(5.1 + row * 0.65)
    badge = add_shape_rect(sl, x, y, Inches(1.55), Inches(0.45), LIGHT_BG)
    add_text_box(sl, x + Inches(0.05), y + Inches(0.05), Inches(1.45), Inches(0.35),
                 section, font_size=10, color=DARK, alignment=PP_ALIGN.CENTER)

add_text_box(sl, Inches(0.8), Inches(6.6), Inches(11), Inches(0.4),
             "Compatible con Microsoft Word y Google Docs. Totalmente editable. Branding con colores y logo de la organización.",
             font_size=12, color=MEDIUM_GRAY)


# ── SLIDE 15: Seguridad ──
sl = prs.slides.add_slide(BLANK)
make_section_header(sl, 14, "Seguridad y privacidad",
                    "Arquitectura diseñada para proteger datos sensibles")

security_data = [
    ["Capa", "Mecanismo", "Detalle"],
    ["Aislamiento multi-tenant", "Row-Level Security (RLS)", "Políticas en PostgreSQL a nivel de fila"],
    ["Separación de PII", "Tablas separadas", "participants (nombre/email) ≠ respondents (anónimo)"],
    ["Acceso a encuesta", "Token único", "Sin email ni login requerido para responder"],
    ["Anonimato estadístico", "Umbral k=5", "Segmentos con <5 personas excluidos automáticamente"],
    ["Control de acceso", "Roles + DB-level", "super_admin, org_admin, member (RLS enforced)"],
    ["Autenticación", "Magic link", "Sin contraseñas — link por email"],
    ["Calidad de datos", "Checks de atención", "2 ítems de validación + ítems invertidos"],
]
add_table(sl, Inches(0.8), Inches(1.5), Inches(11.7), Inches(4.5), security_data,
          col_widths=[Inches(3), Inches(2.5), Inches(6.2)])


# ── SLIDE 16: Stack tecnológico ──
sl = prs.slides.add_slide(BLANK)
make_section_header(sl, 15, "Stack tecnológico",
                    "Tecnología moderna, escalable y mantenible")

tech_data = [
    ["Componente", "Tecnología", "Por qué"],
    ["Framework", "Next.js 16 (App Router)", "Server Components, SSR, serverless"],
    ["Lenguaje", "TypeScript", "Type-safety end-to-end"],
    ["Base de datos", "Supabase (PostgreSQL)", "Auth + RLS + Storage + Realtime"],
    ["UI", "shadcn/ui + Tailwind CSS v4", "Accesible (Radix), personalizable"],
    ["Gráficos", "Recharts", "21 componentes reutilizables"],
    ["Validación", "Zod + React Hook Form", "Schemas compartidos client/server"],
    ["Email", "Resend", "Transaccional, HTML templates"],
    ["IA", "Anthropic / OpenAI / Ollama", "Triple backend con fallback"],
    ["ONA", "Python (igraph + matplotlib)", "Kernel C — alta performance"],
    ["Exportación", "docx + exceljs", "DOCX editable + Excel multi-hoja"],
    ["Deploy", "Vercel", "Edge, serverless, CI/CD automático"],
]
add_table(sl, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.5), tech_data,
          col_widths=[Inches(2), Inches(3.5), Inches(6.2)])


# ── SLIDE 17: Para quién ──
sl = prs.slides.add_slide(BLANK)
make_section_header(sl, 16, "¿Para quién es ClimaLab?", "Mercado objetivo")

# Size
add_shape_rect(sl, Inches(0.8), Inches(1.5), Inches(5.5), Inches(2.0), LIGHT_BG)
add_text_box(sl, Inches(1.1), Inches(1.65), Inches(5), Inches(0.4),
             "Tamaño de empresa", font_size=16, color=GREEN, bold=True)
add_text_box(sl, Inches(1.1), Inches(2.1), Inches(5), Inches(1.0),
             "20 a 500+ colaboradores\nPYMEs y medianas empresas que necesitan datos\npara gestionar su clima, no solo intuición.",
             font_size=14, color=DARK)

# Geography
add_shape_rect(sl, Inches(7), Inches(1.5), Inches(5.8), Inches(2.0), LIGHT_BG)
add_text_box(sl, Inches(7.3), Inches(1.65), Inches(5.2), Inches(0.4),
             "Geografía y contexto", font_size=16, color=GREEN, bold=True)
add_text_box(sl, Inches(7.3), Inches(2.1), Inches(5.2), Inches(1.0),
             "LATAM — diseñado para el contexto cultural\ny lingüístico latinoamericano.\nTodos los sectores (módulos sectoriales en roadmap).",
             font_size=14, color=DARK)

# Decision makers
add_text_box(sl, Inches(0.8), Inches(4.0), Inches(5), Inches(0.4),
             "Decisores", font_size=18, color=DARK, bold=True)

personas = [
    ("Recursos Humanos", "Diagnóstico y seguimiento\nde clima organizacional"),
    ("Gerencia General", "Decisiones basadas en datos\nsobre cultura y gente"),
    ("Transformación", "Medición de impacto\nde iniciativas de cambio"),
    ("Consultores OD", "Herramienta escalable\npara múltiples clientes"),
]
for i, (role, desc) in enumerate(personas):
    x = Inches(0.8 + i * 3.1)
    add_shape_rect(sl, x, Inches(4.5), Inches(2.8), Inches(2.0), WHITE)
    add_shape_rect(sl, x, Inches(4.5), Inches(2.8), Inches(0.06), GREEN)
    add_text_box(sl, x + Inches(0.2), Inches(4.7), Inches(2.4), Inches(0.35),
                 role, font_size=14, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.2), Inches(5.1), Inches(2.4), Inches(1.0),
                 desc, font_size=12, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)


# ── SLIDE 18: Ventajas competitivas ──
sl = prs.slides.add_slide(BLANK)
make_section_header(sl, 17, "Ventajas competitivas", "¿Por qué elegir ClimaLab?")

advantages = [
    ["Ventaja", "Detalle"],
    ["Rigor científico", "22 dimensiones validadas (vs. 6-10 genéricas)"],
    ["Análisis psicométrico", "rwg, Alpha, Pearson — nivel de investigación académica"],
    ["Red perceptual (ONA)", "Único en mercado PYME: comunidades y puentes culturales"],
    ["IA triple backend", "Narrativas automáticas con fallback inteligente"],
    ["Marca blanca completa", "Encuesta, emails, reportes — identidad del cliente"],
    ["Ficha técnica auto-generada", "Limitaciones y metodología transparente"],
    ["LATAM-nativo", "Contexto cultural y lingüístico latinoamericano"],
    ["Exportación editable", "DOCX 14 secciones (no PDF estático)"],
    ["Anonimato real", "PII separado, umbrales de segmento, tokens sin login"],
    ["Multi-instrumento", "Base + 3 módulos opcionales por campaña"],
]
add_table(sl, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.3), advantages,
          col_widths=[Inches(3.5), Inches(8.2)])


# ── SLIDE 19: Roadmap ──
sl = prs.slides.add_slide(BLANK)
make_section_header(sl, 18, "Roadmap — Evolución del producto",
                    "Tres horizontes de desarrollo")

horizons = [
    ("H1: Operacional", "Completado — v4.9", [
        "Core v4.0 + Pulso + 3 módulos opcionales",
        "Dashboard 11 sub-páginas + 21 componentes",
        "IA triple backend (Anthropic + DGX + Ollama)",
        "ONA perceptual (Leiden + NMI)",
        "DOCX/Excel/CSV + marca blanca + emails",
    ], GREEN),
    ("H2: Analítico", "6–18 meses", [
        "Análisis Factorial Confirmatorio (CFA)",
        "Invarianza de medición (cross-org)",
        "Normas regionales LATAM",
        "ONA sociométrica (interacción real)",
        "Intervalos de confianza para segmentos",
    ], CYAN),
    ("H3: Avanzado", "18–36 meses", [
        "Modelado Lineal Jerárquico (HLM)",
        "NLP para clasificación de comentarios",
        "Módulos sectoriales (salud, retail...)",
        "API e integraciones (HRIS, BI, webhooks)",
        "Dashboard de administrador multi-org",
    ], RGBColor(0x2F, 0x5D, 0xFF)),
]

for i, (title, timeline, items, color) in enumerate(horizons):
    x = Inches(0.8 + i * 4.1)
    add_shape_rect(sl, x, Inches(1.5), Inches(3.7), Inches(5.3), WHITE)
    add_shape_rect(sl, x, Inches(1.5), Inches(3.7), Inches(0.06), color)
    add_text_box(sl, x + Inches(0.2), Inches(1.7), Inches(3.3), Inches(0.35),
                 title, font_size=18, color=color, bold=True)
    add_text_box(sl, x + Inches(0.2), Inches(2.1), Inches(3.3), Inches(0.3),
                 timeline, font_size=12, color=MEDIUM_GRAY)
    add_bullet_list(sl, x + Inches(0.2), Inches(2.5), Inches(3.3), Inches(4.0),
                    items, font_size=12, spacing=Pt(6))


# ── SLIDE 20: Flujo de trabajo ──
sl = prs.slides.add_slide(BLANK)
make_section_header(sl, 19, "Flujo de trabajo completo",
                    "Del diagnóstico a la acción en 7 pasos")

workflow = [
    ("1", "Configura", "Crea campaña, selecciona\ninstrumento + módulos"),
    ("2", "Personaliza", "Aplica marca: logo,\ncolores, textos"),
    ("3", "Invita", "Emails de invitación con\nenlace único por participante"),
    ("4", "Mide", "Encuesta anónima (~10 min)\ncon respaldo automático"),
    ("5", "Recuerda", "Recordatorios a quienes\nno han completado"),
    ("6", "Analiza", "Cálculo: estadísticas +\nONA + alertas + IA"),
    ("7", "Actúa", "11 módulos de resultados +\nreporte DOCX editable"),
]

for i, (num, title, desc) in enumerate(workflow):
    x = Inches(0.3 + i * 1.85)
    # Number circle
    circle = sl.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.55), Inches(1.7), Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = GREEN if i < 6 else CYAN
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    if i < 6:
        add_text_box(sl, x + Inches(1.25), Inches(1.75), Inches(0.7), Inches(0.5),
                     "→", font_size=22, color=LIGHT_GRAY, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(sl, x, Inches(2.5), Inches(1.7), Inches(0.35),
                 title, font_size=14, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x, Inches(2.9), Inches(1.7), Inches(1.0),
                 desc, font_size=11, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom highlight
add_shape_rect(sl, Inches(0.8), Inches(4.5), Inches(11.7), Inches(2.5), LIGHT_BG)
add_text_box(sl, Inches(1.2), Inches(4.7), Inches(10), Inches(0.4),
             "Resultados que puedes esperar:", font_size=18, color=DARK, bold=True)

results = [
    "Diagnóstico claro con indicadores de rigor científico",
    "Identificación de perfiles de compromiso en tu organización",
    "Comparativas por área, antigüedad y segmentos",
    "Recomendaciones accionables para mejorar el clima",
    "Mapa de percepción organizacional (ONA)",
    "Narrativas IA que traducen datos en historias claras",
    "Reporte ejecutivo editable listo para presentar a la directiva",
]
add_bullet_list(sl, Inches(1.2), Inches(5.2), Inches(10), Inches(2.0),
                results, font_size=13, spacing=Pt(4))


# ── SLIDE 21: Cierre ──
sl = prs.slides.add_slide(BLANK)
set_slide_bg(sl, GREEN)
add_shape_rect(sl, Inches(0), Inches(0), Inches(13.33), Inches(0.12), RGBColor(0x1E, 0x7A, 0x38))

add_text_box(sl, Inches(0.8), Inches(1.5), Inches(11), Inches(0.6),
             "ClimaLab", font_size=28, color=WHITE, bold=True)

add_text_box(sl, Inches(0.8), Inches(2.5), Inches(11), Inches(1.5),
             "Mide lo que importa\nen tu organización",
             font_size=48, color=WHITE, bold=True)

add_text_box(sl, Inches(0.8), Inches(4.5), Inches(11), Inches(0.8),
             "22 dimensiones · 11 módulos de resultados · IA integrada\nONA perceptual · Marca blanca · Exportación editable",
             font_size=20, color=RGBColor(0xD0, 0xF0, 0xD8))

add_shape_rect(sl, Inches(0), Inches(6.2), Inches(13.33), Inches(1.3), RGBColor(0x1E, 0x7A, 0x38))
add_text_box(sl, Inches(0.8), Inches(6.35), Inches(5), Inches(0.4),
             "climalab.vercel.app", font_size=20, color=WHITE, bold=True)
add_text_box(sl, Inches(0.8), Inches(6.8), Inches(5), Inches(0.4),
             "Un producto de Rizo.ma Consulting  ·  Panamá",
             font_size=14, color=RGBColor(0xA0, 0xD8, 0xB0))

# ── Save ──
output_path = "docs/ClimaLab_Presentacion.pptx"
prs.save(output_path)
print(f"Presentación generada: {output_path}")
print(f"  → {len(prs.slides)} diapositivas")
print(f"  → Formato: 16:9 widescreen")
print(f"  → Colores: Rizoma brand (#289448, #1FACC0, #C32421)")
